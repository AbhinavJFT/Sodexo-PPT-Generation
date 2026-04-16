"""Rewrite a .pptx's text in-place for a new topic, preserving all visuals."""
import json
import logging
from typing import Any, Dict, Iterator, List

from pptx import Presentation
from pptx.util import Pt

log = logging.getLogger("rewrite")

REWRITE_SYSTEM = """You are rewriting the text of a PowerPoint deck so it covers a new topic while the visual structure stays identical.

Output contract:
- Return ONLY JSON: {"rewrites": {"<id>": "<new_text>", ...}}
- Use the exact "id" values you were given. Never invent new ids.
- Every id in the input MUST appear in the output.

Content rules:
- Keep the same character count as the original where possible (hard limit: no more than 1.5x the original length, ideally ±20%).
- Preserve role: title stays a title, bullet stays a bullet-sized fragment, section number stays a number.
- Do NOT rewrite any of these — return the ORIGINAL text verbatim:
  * Boilerplate slides (Instructions for use, Credits, Icon catalogs, Font names, "Alternative resources", "Resources", "Thanks!" credits pages).
  * Decorative numbers ("01", "02", "2024", percentages, currency values).
  * Brand names, placeholder names like "MARS".
  * Strings that are just punctuation or whitespace.
- For real content slides (title slide, introduction, scope, timeline, deliverables, etc.), rewrite everything to fit the new topic with coherent, specific language — not generic filler."""


def _iter_paragraphs(slide) -> Iterator[Dict[str, Any]]:
    """Yield dicts describing every non-empty paragraph in a slide."""
    for shape_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text
                if text.strip():
                    yield {
                        "shape_idx": shape_idx,
                        "para_idx": para_idx,
                        "text": text,
                        "role": "title" if shape == slide.shapes.title else "text",
                    }
        if shape.has_table:
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    for para_idx, para in enumerate(cell.text_frame.paragraphs):
                        text = para.text
                        if text.strip():
                            yield {
                                "shape_idx": shape_idx,
                                "row": row_idx,
                                "col": col_idx,
                                "para_idx": para_idx,
                                "text": text,
                                "role": "table_cell",
                            }


def extract_entries(prs) -> List[Dict[str, Any]]:
    """Collect addressable text entries from every slide."""
    entries = []
    for slide_idx, slide in enumerate(prs.slides):
        for seq, ref in enumerate(_iter_paragraphs(slide)):
            ref["slide"] = slide_idx
            ref["id"] = f"{slide_idx}-{seq}"
            entries.append(ref)
    return entries


def _get_paragraph(prs, entry: Dict[str, Any]):
    slide = prs.slides[entry["slide"]]
    shape = list(slide.shapes)[entry["shape_idx"]]
    if "row" in entry:
        cell = shape.table.rows[entry["row"]].cells[entry["col"]]
        return cell.text_frame.paragraphs[entry["para_idx"]]
    return shape.text_frame.paragraphs[entry["para_idx"]]


def set_paragraph_text(para, new_text: str) -> None:
    """Replace paragraph text using the first run's formatting; drop extra runs."""
    runs = list(para.runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        r = para.add_run()
        r.text = new_text


def apply_rewrites(prs, entries: List[Dict[str, Any]], rewrites: Dict[str, str]) -> int:
    n = 0
    for e in entries:
        new_text = rewrites.get(e["id"])
        if not isinstance(new_text, str):
            continue
        if new_text == e["text"]:
            continue
        try:
            para = _get_paragraph(prs, e)
            set_paragraph_text(para, new_text)
            n += 1
        except Exception:
            log.exception("failed to apply rewrite for %s", e["id"])
    return n


def _chunk(items, size):
    for i in range(0, len(items), size):
        yield items[i : i + size]


def request_rewrites(groq_client, model: str, topic: str, entries: List[Dict[str, Any]]) -> Dict[str, str]:
    """Call Groq (chunking by slide groups if large) and merge the rewrites."""
    # Group entries into chunks of at most ~60 entries, splitting only on slide boundaries.
    chunks: List[List[Dict[str, Any]]] = []
    current: List[Dict[str, Any]] = []
    current_size = 0
    last_slide = -1
    for e in entries:
        if current and e["slide"] != last_slide and current_size >= 60:
            chunks.append(current)
            current = []
            current_size = 0
        current.append(e)
        current_size += 1
        last_slide = e["slide"]
    if current:
        chunks.append(current)

    log.info("rewrite: %d entries across %d chunks", len(entries), len(chunks))
    merged: Dict[str, str] = {}
    for i, chunk in enumerate(chunks):
        payload = [
            {"id": e["id"], "slide": e["slide"], "role": e["role"], "len": len(e["text"]), "text": e["text"]}
            for e in chunk
        ]
        user_msg = (
            f"New topic:\n{topic}\n\n"
            "Rewrite the text below for this topic. Follow every rule in the system prompt. "
            'Return JSON: {"rewrites": {"<id>": "<new_text>", ...}}.\n\n'
            f"Entries:\n{json.dumps(payload, ensure_ascii=False)}"
        )
        resp = groq_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": REWRITE_SYSTEM},
                {"role": "user", "content": user_msg},
            ],
            response_format={"type": "json_object"},
            temperature=0.4,
            max_tokens=8000,
        )
        raw = resp.choices[0].message.content or "{}"
        try:
            data = json.loads(raw)
            rw = data.get("rewrites", {})
            if isinstance(rw, dict):
                for k, v in rw.items():
                    if isinstance(v, str):
                        merged[k] = v
        except json.JSONDecodeError:
            log.warning("chunk %d: bad JSON from groq, skipping", i)
    log.info("rewrite: got %d/%d replacements", len(merged), len(entries))
    return merged


def rewrite_deck(groq_client, model: str, src_path: str, dst_path: str, topic: str) -> Dict[str, Any]:
    prs = Presentation(src_path)
    entries = extract_entries(prs)
    if not entries:
        prs.save(dst_path)
        return {"total_entries": 0, "replaced": 0}
    rewrites = request_rewrites(groq_client, model, topic, entries)
    replaced = apply_rewrites(prs, entries, rewrites)
    prs.save(dst_path)
    return {"total_entries": len(entries), "replaced": replaced, "slides": len(prs.slides)}
