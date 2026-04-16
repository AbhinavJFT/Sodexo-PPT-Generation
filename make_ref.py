"""Create a small sample reference .pptx for testing."""
from pptx import Presentation
from pptx.util import Inches, Pt

p = Presentation()
s = p.slides.add_slide(p.slide_layouts[0])
s.shapes.title.text = "Sample Reference Deck"
s.placeholders[1].text = "Q4 2025 Review"

s2 = p.slides.add_slide(p.slide_layouts[1])
s2.shapes.title.text = "Agenda"
s2.placeholders[1].text = "Overview\nResults\nNext steps"

p.save("/tmp/sample_ref.pptx")
print("wrote /tmp/sample_ref.pptx")
